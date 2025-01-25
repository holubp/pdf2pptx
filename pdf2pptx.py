#!/usr/bin/env python3

import os
import sys
import shutil
import subprocess
import argparse
import logging

from pptx import Presentation
from pptx.util import Inches

def error_and_exit(message, exit_code=1):
    """Log an error message and exit."""
    logging.error(message)
    sys.exit(exit_code)

def parse_slide_size(size_str):
    """
    Parse a string describing the slide size.
    Options:
      - "16:9" -> (13.3333, 7.5)
      - "4:3"  -> (10, 7.5)
      - "WxH"  -> (float(W), float(H)) in inches
    Returns (width_in, height_in) as floats in inches, or exits on error.
    """
    if size_str in ("16:9", "16x9"):
        return (13.3333, 7.5)
    elif size_str in ("4:3", "4x3"):
        return (10.0, 7.5)
    else:
        # Attempt custom "WxH" format in inches
        try:
            w_str, h_str = size_str.split('x')
            return (float(w_str), float(h_str))
        except Exception:
            error_and_exit(
                "Could not parse slide size. Use '16:9', '4:3', or 'WxH' in inches."
            )

def main():
    parser = argparse.ArgumentParser(
        description="Convert a PDF to PPTX (one PDF page per slide) using ImageMagick and python-pptx."
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="show detailed progress information (INFO level logs)"
    )
    parser.add_argument(
        "--retain-image-files",
        action="store_true",
        help="do not delete the temporary image files/directory after creating the PPTX"
    )
    parser.add_argument(
        "-d", "--dpi",
        type=int,
        default=1200,
        help="set DPI (dots per inch) used by ImageMagick convert (default: 1200)"
    )
    parser.add_argument(
        "-q", "--quality",
        type=int,
        default=95,
        help="set image compression quality (1-100) for ImageMagick (default: 95)"
    )
    parser.add_argument(
        "-f", "--image-format",
        choices=["JPEG", "PNG"],
        default="JPEG",
        help="choose intermediate image format: JPEG or PNG (default: JPEG)"
    )
    parser.add_argument(
        "-t", "--temp-dir",
        help="directory to store temporary image files (must not already exist); "
             "defaults to <PDF_basename>_temp"
    )
    parser.add_argument(
        "-s", "--slide-size",
        default="16:9",
        help="slide size: '16:9', '4:3', or 'WxH' in inches (default: 16:9 => 13.3333x7.5)"
    )
    parser.add_argument("input_pdf", help="the input PDF file")
    parser.add_argument("output_pptx", help="the output PPTX file")
    args = parser.parse_args()

    # Configure logging
    # Default to WARNING (only warnings/errors). If verbose, switch to INFO.
    logging.basicConfig(
        level=logging.WARNING,
        format="[%(levelname)s] %(message)s"
    )
    if args.verbose:
        logging.getLogger().setLevel(logging.INFO)

    pdf_file = args.input_pdf
    pptx_file = args.output_pptx

    # 1) Validate the input PDF path
    if not os.path.isfile(pdf_file):
        error_and_exit(f"PDF file '{pdf_file}' does not exist.")

    # 2) Determine temporary directory
    base_name, _ = os.path.splitext(os.path.basename(pdf_file))
    if args.temp_dir:
        temp_dir = args.temp_dir
    else:
        temp_dir = f"{base_name}_temp"

    if os.path.exists(temp_dir):
        error_and_exit(f"Temporary directory '{temp_dir}' already exists. "
                       "Remove it or specify a different directory with --temp-dir.")

    # Attempt to create temp directory
    try:
        os.mkdir(temp_dir)
    except OSError as e:
        error_and_exit(f"Could not create directory '{temp_dir}': {e}")

    # Determine file extension based on requested image format
    if args.image_format == "JPEG":
        extension = "jpg"
    else:  # "PNG"
        extension = "png"

    # 3) Build the ImageMagick 'convert' command
    # Place '-verbose' right after "convert" if needed
    convert_cmd = ["convert"]
    if args.verbose:
        convert_cmd.append("-verbose")

    convert_cmd += [
        "-density", str(args.dpi),
        pdf_file,
        "-resize", "50%",
        "-quality", str(args.quality),
        os.path.join(temp_dir, f"page.{extension}")
    ]

    logging.info("Running ImageMagick convert:\n  %s", " ".join(convert_cmd))
    try:
        subprocess.run(convert_cmd, check=True)
    except subprocess.CalledProcessError as e:
        # Cleanup if conversion fails and user doesn't want to retain images
        if os.path.isdir(temp_dir) and not args.retain_image_files:
            shutil.rmtree(temp_dir, ignore_errors=True)
        error_and_exit(f"ImageMagick conversion failed: {e}")

    logging.info("Conversion done. Now generating PPTX...")

    # 4) Parse slide size from user input (or default)
    slide_width_in, slide_height_in = parse_slide_size(args.slide_size)

    # 5) Create the PPTX
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    blank_slide_layout = prs.slide_layouts[6]

    # 6) Collect image files and add them as slides
    index = 0
    slide_count = 0

    while True:
        image_filename = os.path.join(temp_dir, f"page-{index}.{extension}")
        if not os.path.isfile(image_filename):
            break

        logging.info("Adding slide for '%s'", image_filename)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            image_filename,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        index += 1
        slide_count += 1

    # 7) Save PPTX
    prs.save(pptx_file)
    logging.info("Successfully created '%s' with %d slide(s).", pptx_file, slide_count)

    # 8) Cleanup
    if not args.retain_image_files:
        logging.info("Removing temporary directory '%s'...", temp_dir)
        try:
            shutil.rmtree(temp_dir)
            logging.info("Removed '%s'.", temp_dir)
        except OSError as e:
            logging.warning("Could not remove '%s': %s", temp_dir, e)
    else:
        logging.info("Retaining all image files in '%s'.", temp_dir)

if __name__ == "__main__":
    main()

